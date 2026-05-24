"""Generate SUNUM.pptx — final presentation deck.

15 slides, 13.5 min talk + 1.5 min Q&A buffer. Visual identity mirrors
the midterm deck: dark navy background, amber-orange brand accent,
section badges, underline accent under titles, card-based layouts.

Run from the project root:
    backend/.venv/bin/python slides/build.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Color palette (mirrors midterm visual identity + our brand tokens)
# ---------------------------------------------------------------------------
C = {
    "bg":           RGBColor(0x0A, 0x0F, 0x1E),   # near-black navy
    "bg_glow":      RGBColor(0x14, 0x1D, 0x36),
    "card_bg":      RGBColor(0x13, 0x1A, 0x2E),
    "card_bg2":     RGBColor(0x1A, 0x24, 0x3A),
    "card_border":  RGBColor(0x2A, 0x37, 0x52),
    "text":         RGBColor(0xF1, 0xF5, 0xF9),
    "text_dim":     RGBColor(0x94, 0xA3, 0xB8),
    "text_subtle":  RGBColor(0x64, 0x74, 0x8B),
    "amber":        RGBColor(0xFB, 0xBF, 0x24),
    "fire":         RGBColor(0xF9, 0x73, 0x16),
    "fire_strong":  RGBColor(0xEA, 0x58, 0x0C),
    "teal":         RGBColor(0x14, 0xB8, 0xA6),
    "blue":         RGBColor(0x60, 0xA5, 0xFA),
    "purple":       RGBColor(0x8B, 0x5C, 0xF6),
    "red":          RGBColor(0xEF, 0x44, 0x44),
    "rose":         RGBColor(0xF4, 0x3F, 0x5E),
    "green":        RGBColor(0x10, 0xB9, 0x81),
    "yellow":       RGBColor(0xFA, 0xCC, 0x15),
}

FONT = "Inter"            # falls back gracefully if not installed
FONT_MONO = "JetBrains Mono"
N_SLIDES = 15


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def color_of(c):
    return C[c] if isinstance(c, str) else c


def add_slide(prs: Presentation):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C["bg"]
    return slide


def add_text(slide, left, top, width, height, lines, *,
             size=14, color="text", bold=False, italic=False,
             align="left", anchor="top", font=FONT, line_spacing=1.15):
    """Add a text box. `lines` is either a string (single line) or a list of
    dicts {text, size?, color?, bold?, italic?, font?} for mixed runs/paragraphs.

    Lists of strings get rendered as paragraphs with the default style.
    """
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    if anchor == "middle": tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom": tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    if isinstance(lines, str):
        lines = [{"text": lines}]
    if isinstance(lines, list) and lines and isinstance(lines[0], str):
        lines = [{"text": s} for s in lines]

    for i, item in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.line_spacing = line_spacing
        if align == "center": p.alignment = PP_ALIGN.CENTER
        elif align == "right": p.alignment = PP_ALIGN.RIGHT

        if "runs" in item:
            for j, r in enumerate(item["runs"]):
                run = p.add_run() if j > 0 else (p.add_run() if not p.runs else p.runs[0])
                # ensure run exists
                if not p.runs:
                    run = p.add_run()
                else:
                    run = p.runs[j] if j < len(p.runs) else p.add_run()
                run.text = r["text"]
                run.font.size = Pt(r.get("size", size))
                run.font.color.rgb = color_of(r.get("color", color))
                run.font.bold = r.get("bold", bold)
                run.font.italic = r.get("italic", italic)
                run.font.name = r.get("font", font)
            continue

        run = p.add_run()
        run.text = item.get("text", "")
        run.font.size = Pt(item.get("size", size))
        run.font.color.rgb = color_of(item.get("color", color))
        run.font.bold = item.get("bold", bold)
        run.font.italic = item.get("italic", italic)
        run.font.name = item.get("font", font)
    return tb


def add_rect(slide, left, top, width, height, *,
             fill="card_bg", line=None, line_width=0.75,
             rounded=True, corner=0.04):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color_of(fill)
    if line is not None:
        shape.line.color.rgb = color_of(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    if rounded:
        # tweak corner radius (default is too round)
        try:
            shape.adjustments[0] = corner
        except Exception:
            pass
    # remove default shadow
    sp = shape.shadow
    try:
        sp.inherit = False
    except Exception:
        pass
    return shape


def add_line(slide, x1, y1, x2, y2, *, color="amber", width=2.5):
    shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = color_of(color)
    shape.line.width = Pt(width)
    return shape


def add_section_badge(slide, label, color="amber", left=0.55, top=0.45):
    """Small rounded chip in the upper-left of a slide, like midterm."""
    rgb = color_of(color)
    # convert RGB to a soft fill (alpha simulated by mixing with bg)
    fill = RGBColor(
        int(0.18 * rgb[0] + 0.82 * C["bg"][0]),
        int(0.18 * rgb[1] + 0.82 * C["bg"][1]),
        int(0.18 * rgb[2] + 0.82 * C["bg"][2]),
    )
    w = max(1.0, 0.085 * len(label) + 0.45)
    rect = add_rect(slide, left, top, w, 0.34, fill=fill, line=color, line_width=0.6, rounded=True, corner=0.45)
    add_text(slide, left + 0.05, top + 0.03, w - 0.1, 0.28, label,
             size=10.5, color=color, bold=True, align="center", anchor="middle")
    return rect


def add_title(slide, text, *, top=0.95, left=0.55, width=10, color="text", size=36):
    add_text(slide, left, top, width, 0.85, text,
             size=size, color=color, bold=True, font=FONT)
    add_line(slide, left, top + 0.78, left + 1.0, top + 0.78,
             color="amber", width=3.0)


def add_page_number(slide, idx, total=N_SLIDES, top=7.18, color="text_subtle"):
    add_text(slide, 12.4, top, 0.85, 0.22, f"{idx} / {total}",
             size=10, color=color, align="right", font=FONT_MONO)


def add_card(slide, left, top, width, height, *, fill="card_bg",
             accent_color=None, accent_side="top"):
    """Card with optional colored accent stripe along one side."""
    card = add_rect(slide, left, top, width, height,
                    fill=fill, line="card_border", line_width=0.6,
                    rounded=True, corner=0.03)
    if accent_color is not None:
        if accent_side == "top":
            add_rect(slide, left + 0.08, top - 0.02, width - 0.16, 0.05,
                     fill=accent_color, rounded=False)
        elif accent_side == "left":
            add_rect(slide, left - 0.02, top + 0.08, 0.05, height - 0.16,
                     fill=accent_color, rounded=False)
    return card


def add_bullet(slide, left, top, width, items, *, size=14, gap=0.32, color="text", bullet_color="amber"):
    """Vertically stacked bullet items."""
    for i, txt in enumerate(items):
        y = top + i * gap
        # bullet dot
        add_rect(slide, left, y + 0.10, 0.10, 0.10, fill=bullet_color, rounded=True, corner=0.5)
        add_text(slide, left + 0.22, y, width - 0.22, gap, txt,
                 size=size, color=color)


def add_footer_date(slide, text="May 2026", left=0.6, top=7.18, color="text_subtle"):
    add_text(slide, left, top, 4, 0.22, text, size=10, color=color)


# ===========================================================================
# SLIDE CONTENT
# ===========================================================================

def build_s1_title(prs):
    """Title slide — flame logo, project title, team, chips."""
    s = add_slide(prs)

    # subtle glow behind the flame
    add_rect(s, 9.3, 0.3, 4.0, 4.0,
             fill=RGBColor(0x2A, 0x18, 0x10), rounded=True, corner=0.5)

    # flame "icon" — emoji-styled big text
    add_text(s, 9.2, 0.6, 4.0, 3.6, "🔥",
             size=180, color="amber", align="center", anchor="middle")

    # accent bar above title
    add_rect(s, 0.7, 1.55, 0.55, 0.10, fill="amber", rounded=False)

    # title
    add_text(s, 0.7, 1.75, 8.5, 1.5, "Stop the Fire!",
             size=66, color="text", bold=True)

    # subtitle
    add_text(s, 0.7, 3.0, 8.5, 0.55,
             "A Network Optimization Approach to Wildfire Containment",
             size=22, color="text_dim")

    # meta line
    add_text(s, 0.7, 3.7, 8.5, 0.4,
             "IE 492 — Senior Design Project   ·   Final Presentation",
             size=14, color="text_subtle")

    # team box
    team_card = add_rect(s, 0.7, 4.6, 6.8, 1.85,
                         fill="card_bg", line="card_border", rounded=True, corner=0.03)
    # left amber bar inside card
    add_rect(s, 0.66, 4.74, 0.06, 1.57, fill="amber", rounded=False)

    add_text(s, 0.95, 4.78, 3.0, 0.3, "TEAM MEMBERS",
             size=11, color="amber", bold=True)
    add_text(s, 4.5, 4.78, 3.0, 0.3, "SUPERVISOR: ",
             size=11, color="amber", bold=True)
    add_text(s, 5.45, 4.78, 2.0, 0.3, "Tınaz Ekim",
             size=11, color="text")

    members = [
        "Efe Ergen — 2021402168",
        "Emre Barutçu — 2023402219",
        "Mehmet Efe Aloğlu — 2020402024",
        "Kerem Külünkoğlu — 2023402222",
    ]
    for i, m in enumerate(members):
        add_text(s, 0.95, 5.18 + i * 0.30, 6.0, 0.28, m,
                 size=13, color="text")

    # chips on the right
    chips = [
        ("Real-world graph", "teal"),
        ("13 strategies",    "purple"),
        ("Decision support", "fire"),
    ]
    chip_w = 1.95
    for i, (label, col) in enumerate(chips):
        rgb = color_of(col)
        fill_rgb = RGBColor(
            int(0.20 * rgb[0] + 0.80 * C["bg"][0]),
            int(0.20 * rgb[1] + 0.80 * C["bg"][1]),
            int(0.20 * rgb[2] + 0.80 * C["bg"][2]),
        )
        x = 8.5 + (i % 2) * (chip_w + 0.15)
        y = 5.05 + (i // 2) * 0.55
        if i == 2:  # third chip below
            x = 8.5
            y = 5.6
        add_rect(s, x, y, chip_w, 0.42,
                 fill=fill_rgb, line=col, line_width=0.6, rounded=True, corner=0.4)
        add_text(s, x, y, chip_w, 0.42, label,
                 size=11, color=col, bold=True, align="center", anchor="middle")

    add_footer_date(s)
    add_page_number(s, 1)


def build_s2_motivation(prs):
    s = add_slide(prs)
    add_section_badge(s, "MOTIVATION", color="fire")
    add_title(s, "Why now")

    # Big single focus card — Reactive → Proactive
    card = add_card(s, 2.5, 2.4, 8.3, 2.8, fill="card_bg",
                    accent_color="purple", accent_side="top")

    add_text(s, 2.8, 2.65, 7.7, 0.6, "💡  Reactive → Proactive",
             size=30, color="purple", bold=True)
    add_text(s, 2.8, 3.35, 7.7, 1.7,
             "Current wildfire response is reactive — we propose strategic, "
             "graph-theoretic decision-making before the front hits a village. "
             "Stop the Fire! is a concrete prototype for the Turkish "
             "Mediterranean belt — the Köyceğiz–Marmaris region.",
             size=15, color="text", line_spacing=1.4)

    # supporting points (verbal in talk, present as small chips at bottom)
    supports = [
        ("🌲", "Environmental",  "teal"),
        ("🏠", "Human cost",     "red"),
        ("$",  "Economic damage", "amber"),
    ]
    for i, (icon, label, col) in enumerate(supports):
        x = 2.5 + i * 2.8
        add_rect(s, x, 5.7, 2.55, 0.55,
                 fill="card_bg2", line="card_border", line_width=0.5, rounded=True, corner=0.15)
        add_text(s, x + 0.15, 5.78, 0.5, 0.4, icon,
                 size=18, color=col, anchor="middle")
        add_text(s, x + 0.7, 5.78, 1.8, 0.4, label,
                 size=12, color="text_dim", anchor="middle")

    add_text(s, 0.6, 6.6, 12.0, 0.35,
             "Verbal: environmental, human, economic context — covered briefly in narration.",
             size=10, color="text_subtle", italic=True)

    add_page_number(s, 2)


def build_s3_problem(prs):
    s = add_slide(prs)
    add_section_badge(s, "CRITERION 1: PROBLEM", color="red")
    add_title(s, "The Firefighter Problem")

    # Left: formal statement card
    fs_card = add_card(s, 0.6, 2.4, 6.8, 1.7, fill="card_bg",
                       accent_color="amber", accent_side="left")
    add_text(s, 0.85, 2.55, 6.3, 0.4, "FORMAL STATEMENT",
             size=11, color="amber", bold=True)
    add_text(s, 0.85, 2.95, 6.3, 1.1,
             "\"Given a graph G, a fire origin, and k resources per turn, "
             "minimize total burned vertices.\"",
             size=15, color="text", italic=True, line_spacing=1.35)

    # Bullets below
    bullets = [
        "Classical combinatorial optimization · NP-Hard",
        "Deterministic network optimization with limited intervention",
        "Our adaptation: planar / map-based graphs on real geography",
    ]
    add_bullet(s, 0.7, 4.4, 6.7, bullets, size=14, gap=0.42,
               bullet_color="amber")

    # Right: simplified toy graph illustration (boxes representing vertices)
    add_text(s, 8.4, 2.2, 4.4, 0.3, "TOY EXAMPLE",
             size=10, color="text_subtle", bold=True, align="center")

    # build a tiny network visual (3 layers, simulated)
    nodes = [
        (10.6, 2.7, "red",   "A", "Burned"),
        (9.4,  3.5, "red",   "B", "Burned"),
        (11.7, 3.5, "red",   "C", "Burned"),
        (8.5,  4.4, "text_dim",  "D", "Safe"),
        (10.0, 4.4, "green", "E", "Protected"),
        (11.5, 4.4, "green", "F", "Protected"),
        (12.8, 4.4, "red",   "G", "Burned"),
        (9.2,  5.4, "text_dim",  "H", "Safe"),
        (10.5, 5.4, "text_dim",  "I", "Safe"),
        (11.8, 5.4, "text_dim",  "J", "Safe"),
    ]
    # edges (subset for clarity)
    edges = [
        (0,1),(0,2),(1,3),(1,4),(2,5),(2,6),(3,7),(4,7),(4,8),(5,8),(5,9),(6,9),
    ]
    for a, b in edges:
        x1, y1 = nodes[a][0] + 0.18, nodes[a][1] + 0.18
        x2, y2 = nodes[b][0] + 0.18, nodes[b][1] + 0.18
        ln = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        ln.line.color.rgb = C["card_border"]
        ln.line.width = Pt(0.9)

    for (x, y, col, label, _state) in nodes:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.36), Inches(0.36))
        c.fill.solid()
        c.fill.fore_color.rgb = color_of(col)
        c.line.color.rgb = C["bg"]
        c.line.width = Pt(1.2)
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(0.36), Inches(0.36))
        tf = tb.text_frame
        tf.margin_left = Pt(0); tf.margin_right = Pt(0)
        tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(10); r.font.bold = True; r.font.name = FONT
        r.font.color.rgb = C["text"] if col != "text_dim" else C["text_subtle"]

    # legend
    legend_items = [("red", "Burned"), ("green", "Protected"), ("text_dim", "Safe")]
    for i, (col, lbl) in enumerate(legend_items):
        x = 8.6 + i * 1.55
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(6.1), Inches(0.18), Inches(0.18))
        c.fill.solid(); c.fill.fore_color.rgb = color_of(col); c.line.fill.background()
        add_text(s, x + 0.24, 6.07, 1.1, 0.24, lbl, size=11, color="text_dim")

    add_text(s, 0.6, 6.55, 12, 0.3,
             "Foundation — all 13 strategies share this rule; they differ in which vertex to protect.",
             size=11, color="text_subtle", italic=True)

    add_page_number(s, 3)


def build_s4_evolution(prs):
    """DESIGN EVOLUTION ANCHOR."""
    s = add_slide(prs)
    add_section_badge(s, "DESIGN EVOLUTION", color="amber")
    add_title(s, "From midterm to now")

    # stage-setter line
    add_text(s, 0.6, 1.95, 12.2, 0.4,
             "\"We made promises at midterm — here's what we delivered, across five axes.\"",
             size=14, color="amber", italic=True)

    # comparison table — 6 rows
    rows = [
        ("Axis",             "Midterm (April)",          "Final (May)",                                       "header"),
        ("Strategies",       "4",                        "13 (+ 3 lookahead, Articulation, Betweenness, …)",  None),
        ("Network",          "n=30 synthetic random",    "Real Köyceğiz map (ESA WorldCover + Lloyd's)",      None),
        ("Hybrid",           "\"In progress\"",           "hybrid_density_aware + 3 lookahead variants",        None),
        ("Testing",          "30v × 15 runs",            "4,752 sims + real-map runs",                        None),
        ("Deliverable",      "sim engine",               "Web suite (FastAPI + React)",                       None),
        ("Supervisor critique", "—",                      "\"Min-Cut is the wrong question\" → min_damage_cut (+23%)", "highlight"),
    ]
    table_top = 2.55
    col_lefts = [0.6, 3.3, 7.8]
    col_widths = [2.6, 4.4, 5.2]
    row_h = 0.52
    for i, row in enumerate(rows):
        kind = row[3]
        y = table_top + i * row_h
        if kind == "header":
            add_rect(s, col_lefts[0], y, sum(col_widths), row_h,
                     fill="card_bg2", rounded=True, corner=0.05)
            for j, cell in enumerate(row[:3]):
                add_text(s, col_lefts[j] + 0.15, y + 0.05, col_widths[j] - 0.2, row_h - 0.1,
                         cell, size=11, color="amber", bold=True, anchor="middle")
            continue
        if kind == "highlight":
            add_rect(s, col_lefts[0] - 0.04, y - 0.02, sum(col_widths) + 0.08, row_h + 0.04,
                     fill="card_bg", line="fire", line_width=1.8, rounded=True, corner=0.04)
            color_v = ["text", "text", "fire"]
            bolds  = [True, False, True]
        else:
            color_v = ["text_dim", "text", "text"]
            bolds  = [True, False, False]
        for j, cell in enumerate(row[:3]):
            add_text(s, col_lefts[j] + 0.15, y + 0.05, col_widths[j] - 0.2, row_h - 0.1,
                     cell, size=12, color=color_v[j], bold=bolds[j], anchor="middle")

    add_text(s, 0.6, 6.75, 12.0, 0.3,
             "Bottom row (supervisor critique) bridges to S10 — the design-experience anchor.",
             size=10, color="text_subtle", italic=True)

    add_page_number(s, 4)


def build_s5_pipeline(prs):
    s = add_slide(prs)
    add_section_badge(s, "SYSTEM", color="teal")
    add_title(s, "System pipeline")

    # 4 boxes with arrows
    boxes = [
        ("🌍",  "Real Map",         "bbox + ESA WorldCover\n10m raster",      "teal"),
        ("🕸",   "Density-Aware Graph", "Lloyd's k-means\narea_ha per vertex", "blue"),
        ("⚔",   "13 Strategies",       "Greedy / Structural /\nLookahead",     "purple"),
        ("📋",  "Decision Support",     "Recommendation +\nWeb Suite",          "fire"),
    ]
    box_w = 2.7; box_h = 2.6; gap = 0.45
    total_w = 4 * box_w + 3 * gap
    start_x = (13.333 - total_w) / 2
    y = 3.2
    for i, (icon, title, sub, col) in enumerate(boxes):
        x = start_x + i * (box_w + gap)
        rgb = color_of(col)
        fill_rgb = RGBColor(
            int(0.10 * rgb[0] + 0.90 * C["bg"][0]),
            int(0.10 * rgb[1] + 0.90 * C["bg"][1]),
            int(0.10 * rgb[2] + 0.90 * C["bg"][2]),
        )
        add_rect(s, x, y, box_w, box_h, fill=fill_rgb, line=col, line_width=1.0, rounded=True, corner=0.07)
        # icon
        add_text(s, x, y + 0.2, box_w, 0.7, icon,
                 size=40, color=col, align="center", anchor="middle")
        # title
        add_text(s, x, y + 1.05, box_w, 0.55, title,
                 size=17, color=col, bold=True, align="center", anchor="middle")
        # sub
        add_text(s, x + 0.2, y + 1.65, box_w - 0.4, 0.9, sub,
                 size=12, color="text", align="center", anchor="middle", line_spacing=1.3)
        # arrow
        if i < 3:
            ax1 = x + box_w + 0.05
            ax2 = x + box_w + gap - 0.05
            add_line(s, ax1, y + box_h / 2, ax2, y + box_h / 2,
                     color="amber", width=2.0)
            # arrowhead (triangle)
            t = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(ax2 - 0.15), Inches(y + box_h / 2 - 0.10),
                                   Inches(0.20), Inches(0.20))
            t.fill.solid(); t.fill.fore_color.rgb = C["amber"]; t.line.fill.background()

    add_text(s, 0.6, 6.55, 12, 0.35,
             "End-to-end — bbox in, decision out. Each box gets its own slide next.",
             size=12, color="text_dim", italic=True, align="center")

    add_page_number(s, 5)


def build_s6_map_to_graph(prs):
    s = add_slide(prs)
    add_section_badge(s, "MAP → GRAPH", color="blue")
    add_title(s, "Density-aware graph generation (v2)")

    # left: 3 numbered steps
    steps = [
        ("1", "Forest mask",
         "ESA WorldCover classes 10 (forest) + 20 (shrubland); a coarse cell must be ≥ 55% forest fraction.",
         "teal"),
        ("2", "Lloyd's k-means",
         "Density-weighted, 14 iterations, k-means++ initialization.",
         "blue"),
        ("3", "Snap to densest cell",
         "Centroids lock to the densest eligible cell — no holes in settlements inside forest patches.",
         "purple"),
    ]
    y0 = 2.4
    for i, (num, title, body, col) in enumerate(steps):
        y = y0 + i * 1.30
        # number circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(0.7), Inches(y), Inches(0.55), Inches(0.55))
        circ.fill.solid(); circ.fill.fore_color.rgb = color_of(col); circ.line.fill.background()
        add_text(s, 0.7, y, 0.55, 0.55, num,
                 size=20, color="text", bold=True, align="center", anchor="middle")
        # title
        add_text(s, 1.42, y, 5.5, 0.4, title,
                 size=16, color="text", bold=True)
        # body
        add_text(s, 1.42, y + 0.42, 5.5, 0.85, body,
                 size=11.5, color="text_dim", line_spacing=1.4)

    # right: side-by-side mockups (placeholders)
    cap_y = 2.4
    add_text(s, 7.2, cap_y, 5.7, 0.3, "v1 (old)                           v2 (new)",
             size=11, color="text_subtle", bold=True)
    # two mock graphs as cards
    add_card(s, 7.2, cap_y + 0.35, 2.8, 3.6, fill="card_bg")
    add_card(s, 10.1, cap_y + 0.35, 2.8, 3.6, fill="card_bg")
    # scatter "dots" — placeholder pattern
    import random
    random.seed(2)
    # v1: dots everywhere including bare
    for _ in range(35):
        x = 7.35 + random.random() * 2.5
        y = cap_y + 0.5 + random.random() * 3.3
        col = "fire" if random.random() < 0.3 else "green"
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.12), Inches(0.12))
        d.fill.solid(); d.fill.fore_color.rgb = color_of(col); d.line.fill.background()
    # v2: dots clustered in "forest" (random in restricted region with density bias)
    random.seed(7)
    for _ in range(35):
        # bias toward center-right cluster
        cx, cy = 11.5, cap_y + 2.0
        r = random.random() ** 1.7 * 1.1
        ang = random.random() * 6.283
        x = cx + r * (1.2 * (random.random() - 0.5) + 0.8) - 0.4
        y = cy + r * (1.2 * (random.random() - 0.5) + 0.8)
        x = max(10.2, min(12.85, x))
        y = max(cap_y + 0.45, min(cap_y + 3.85, y))
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.12), Inches(0.12))
        d.fill.solid(); d.fill.fore_color.rgb = C["green"]; d.line.fill.background()

    add_text(s, 7.2, 6.18, 5.7, 0.35,
             "→ Vertices snap into the dense forest core; bare/built terrain no longer hosts any vertex.",
             size=10.5, color="text_subtle", italic=True, align="center")

    # quick stat
    add_card(s, 0.6, 6.3, 6.6, 0.65, fill="card_bg2", accent_color="amber", accent_side="left")
    add_text(s, 0.85, 6.4, 6.4, 0.45,
             "Köyceğiz: 500 ha/vertex → 101 vertices, density 0.86–1.00",
             size=12, color="text", bold=True, anchor="middle")

    add_page_number(s, 6)


def build_s7_edge_filter(prs):
    s = add_slide(prs)
    add_section_badge(s, "MAP → GRAPH", color="blue")
    add_title(s, "Edge filter — stricter than midterm")

    # 5-rule table
    rules = [
        ("nonfuel-gap ≥ 60 m continuous",  "56", "red"),
        ("settlement buffer (600 m)",       "27", "amber"),
        ("river buffer (80 m)",             "20", "blue"),
        ("water class (any pixel)",          "3", "teal"),
        ("avg fuel weight < 1.5",            "—", "purple"),
    ]
    # header row
    add_rect(s, 0.6, 2.4, 12.0, 0.5, fill="card_bg2", rounded=True, corner=0.05)
    add_text(s, 0.85, 2.43, 8.5, 0.5, "Rule",
             size=12, color="amber", bold=True, anchor="middle")
    add_text(s, 9.4,  2.43, 3.0, 0.5, "Köyceğiz blocked count",
             size=12, color="amber", bold=True, anchor="middle", align="center")

    for i, (rule, count, col) in enumerate(rules):
        y = 3.0 + i * 0.55
        add_rect(s, 0.6, y, 12.0, 0.5, fill="card_bg", line="card_border", line_width=0.4, rounded=True, corner=0.05)
        # left color stripe
        add_rect(s, 0.6, y, 0.1, 0.5, fill=col, rounded=False)
        add_text(s, 0.85, y + 0.03, 8.5, 0.45, rule,
                 size=13.5, color="text", anchor="middle")
        add_text(s, 9.4,  y + 0.03, 3.0, 0.45, count,
                 size=14, color=col if count != "—" else "text_subtle",
                 bold=True, anchor="middle", align="center", font=FONT_MONO)

    # footer summary
    add_card(s, 0.6, 6.1, 12.0, 0.7, fill="card_bg2", accent_color="amber", accent_side="left")
    add_text(s, 0.85, 6.2, 11.5, 0.5,
             "5 filters: physical untraversability · hydrography · settlements · bare strips · mean fuel.",
             size=13, color="text", italic=True, anchor="middle")

    add_page_number(s, 7)


def build_s8_portfolio(prs):
    s = add_slide(prs)
    add_section_badge(s, "STRATEGIES", color="purple")
    add_title(s, "Strategy portfolio — 3 philosophies")

    # 3 columns
    cols = [
        ("Greedy", "Local / fast", "blue",
         [("max_degree", False),
          ("max_white_neighbors", True),
          ("articulation_priority", False),
          ("saved_component", False)],
         "\"What's good now?\""),
        ("Structural", "Global topology", "purple",
         [("min_cut_edge_front", False),
          ("min_cut_vertex_front", False),
          ("min_damage_cut", True),
          ("betweenness_front", False)],
         "\"Where's the bottleneck?\""),
        ("Lookahead", "Forward-looking", "rose",
         [("one_step_lookahead", True),
          ("full_rollout_pairs", False),
          ("deep_lookahead_3", False),
          ("hybrid_density_aware", False),
          ("random (baseline)", False)],
         "\"What happens 2 turns later?\""),
    ]

    col_w = 4.05; col_h = 4.5; col_gap = 0.15
    start_x = (13.333 - 3 * col_w - 2 * col_gap) / 2
    y_top = 2.3
    for i, (philo, sub, col, items, motto) in enumerate(cols):
        x = start_x + i * (col_w + col_gap)
        rgb = color_of(col)
        fill_rgb = RGBColor(
            int(0.10 * rgb[0] + 0.90 * C["bg"][0]),
            int(0.10 * rgb[1] + 0.90 * C["bg"][1]),
            int(0.10 * rgb[2] + 0.90 * C["bg"][2]),
        )
        add_rect(s, x, y_top, col_w, col_h, fill=fill_rgb, line=col, line_width=1.0, rounded=True, corner=0.04)
        # header
        add_text(s, x + 0.2, y_top + 0.15, col_w - 0.4, 0.45, philo,
                 size=22, color=col, bold=True)
        add_text(s, x + 0.2, y_top + 0.65, col_w - 0.4, 0.3, sub.upper(),
                 size=10, color="text_subtle", bold=True)
        # divider
        add_line(s, x + 0.2, y_top + 1.05, x + col_w - 0.2, y_top + 1.05,
                 color=col, width=0.8)
        # items
        for j, (name, starred) in enumerate(items):
            yy = y_top + 1.2 + j * 0.42
            if starred:
                add_text(s, x + 0.2, yy, 0.3, 0.32, "★",
                         size=14, color="amber", bold=True)
                add_text(s, x + 0.55, yy, col_w - 0.7, 0.32, name,
                         size=13, color="text", bold=True, font=FONT_MONO)
            else:
                add_text(s, x + 0.55, yy, col_w - 0.7, 0.32, name,
                         size=12, color="text_dim", font=FONT_MONO)
        # motto at bottom of col
        add_text(s, x + 0.2, y_top + col_h - 0.55, col_w - 0.4, 0.4, motto,
                 size=10.5, color=col, italic=True, align="center", anchor="middle")

    # footer note
    add_card(s, 0.6, 6.95, 12.2, 0.55, fill="card_bg2", accent_color="amber", accent_side="left")
    add_text(s, 0.85, 7.02, 11.7, 0.4,
             "★ = best representative.  13 strategies = 11 documented in ALGORITHMS.md §1–§12 + FullRolloutPairs + DeepLookahead3.",
             size=11, color="text_dim", anchor="middle")

    add_page_number(s, 8)


def build_s9_lookahead(prs):
    s = add_slide(prs)
    add_section_badge(s, "STRATEGIES", color="purple")
    add_title(s, "Lookahead breakthrough — one_step_lookahead")

    # left: bullets
    add_text(s, 0.6, 2.3, 6.0, 0.35, "THE IDEA",
             size=11, color="amber", bold=True)
    add_text(s, 0.6, 2.7, 6.0, 0.6,
             "Don't just ask \"what's good now?\".",
             size=20, color="text", bold=True)
    add_text(s, 0.6, 3.35, 6.0, 0.6,
             "Ask \"what happens 2 turns later?\".",
             size=20, color="rose", bold=True)

    bullets = [
        {"text": "Top-6 front candidates × ", "size": 14, "color": "text"},
    ]
    add_text(s, 0.6, 4.2, 6.5, 0.4,
             [{"runs": [
                 {"text": "Top-6 candidates × ", "size": 14, "color": "text"},
                 {"text": "C(6,2) = 15 pairs", "size": 14, "color": "amber", "bold": True},
             ]}], size=14)
    add_text(s, 0.6, 4.6, 6.5, 0.4,
             [{"runs": [
                 {"text": "Per pair simulate: ", "size": 14, "color": "text"},
                 {"text": "protect → spread → greedy → spread", "size": 14, "color": "amber", "bold": True, "font": FONT_MONO},
             ]}])
    add_text(s, 0.6, 5.0, 6.5, 0.4,
             [{"runs": [
                 {"text": "Pick the pair with ", "size": 14, "color": "text"},
                 {"text": "lowest burned count", "size": 14, "color": "amber", "bold": True},
             ]}])

    # Result chip
    add_card(s, 0.6, 5.6, 6.4, 0.95, fill="card_bg2", accent_color="rose", accent_side="left")
    add_text(s, 0.85, 5.7, 6.0, 0.4, "RESULT",
             size=10, color="rose", bold=True)
    add_text(s, 0.85, 6.0, 6.0, 0.5,
             "0.5 ms · 6% relative gain over greedy · #1 overall",
             size=14, color="text", bold=True)

    # right: schematic — candidate pair → 2 steps
    add_text(s, 7.4, 2.2, 5.4, 0.3, "PAIR (v1, v2) SIMULATION",
             size=10, color="text_subtle", bold=True, align="center")

    sch_x = 7.4; sch_y = 2.6; sch_w = 5.5
    steps = ["Now\n(t=0)", "Protect\n(v1, v2)", "Spread\n(t=1)", "Greedy\n(k=2)", "Final\n(t=2)"]
    step_w = (sch_w - 0.4) / len(steps)
    for i, label in enumerate(steps):
        x = sch_x + i * step_w + 0.05
        add_card(s, x, sch_y, step_w - 0.1, 1.1, fill="card_bg",
                 accent_color="rose" if i in (1, 3) else None, accent_side="top")
        add_text(s, x, sch_y + 0.15, step_w - 0.1, 0.9, label,
                 size=11, color="text", bold=False, align="center", anchor="middle", line_spacing=1.2)
        # arrow between
        if i < len(steps) - 1:
            ax = x + step_w - 0.1
            ay = sch_y + 0.55
            t = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(ax), Inches(ay), Inches(0.15), Inches(0.16))
            t.fill.solid(); t.fill.fore_color.rgb = C["amber"]; t.line.fill.background()

    # result chip on right
    add_card(s, 7.4, 4.0, 5.5, 0.85, fill="card_bg2", accent_color="amber", accent_side="left")
    add_text(s, 7.65, 4.1, 5.0, 0.3, "→  REPEAT FOR ALL 15 PAIRS",
             size=11, color="amber", bold=True)
    add_text(s, 7.65, 4.4, 5.0, 0.4,
             "Pick the pair (v1*, v2*) that minimizes burned at t=2.",
             size=12.5, color="text")

    # cost / coverage chips
    chips = [("0.5 ms", "amber"), ("15 sims", "blue"), ("k=2 native", "teal")]
    for i, (label, col) in enumerate(chips):
        x = 7.4 + i * 1.85
        add_rect(s, x, 5.2, 1.65, 0.5, fill="card_bg", line=col, line_width=0.6, rounded=True, corner=0.3)
        add_text(s, x, 5.2, 1.65, 0.5, label,
                 size=12, color=col, bold=True, align="center", anchor="middle", font=FONT_MONO)

    add_text(s, 7.4, 5.95, 5.5, 0.5,
             "The engine simulates the engine — full k=2 protect-then-spread dynamics inside each candidate evaluation.",
             size=11, color="text_dim", italic=True, line_spacing=1.3)

    add_page_number(s, 9)


def build_s10_critique(prs):
    """ANCHOR slide — supervisor critique → reformulation."""
    s = add_slide(prs)
    add_section_badge(s, "DESIGN EXPERIENCE", color="amber")
    add_title(s, "Response to our supervisor's critique")

    # Top half: quote
    quote_card = add_card(s, 0.6, 2.2, 12.2, 1.95, fill="card_bg2",
                          accent_color="red", accent_side="left")
    add_text(s, 0.95, 2.35, 5.0, 0.35, "APRIL 14 — SUPERVISOR",
             size=11, color="red", bold=True)
    add_text(s, 0.95, 2.7, 11.4, 0.6,
             "\" min_cut_vertex_front is asking the wrong question.\"",
             size=20, color="text", bold=True, italic=True)
    add_text(s, 0.95, 3.4, 11.4, 0.7,
             "\"It says 'don't let fire reach the distant vertex,' but our objective is 'minimize total burned.'\"",
             size=14, color="text_dim", italic=True, line_spacing=1.4)

    # arrow down
    arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                               Inches(6.5), Inches(4.25), Inches(0.35), Inches(0.35))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = C["amber"]; arrow.line.fill.background()

    # Bottom half: response
    resp = add_card(s, 0.6, 4.7, 12.2, 2.05, fill="card_bg",
                    accent_color="green", accent_side="left")
    add_text(s, 0.95, 4.82, 5.0, 0.35, "OUR REFORMULATION  —  min_damage_cut",
             size=11, color="green", bold=True)
    add_text(s, 0.95, 5.17, 11.4, 0.45,
             "score = saved / cut_size   (ratio, not absolute)",
             size=18, color="text", bold=True, font=FONT_MONO)

    # 3 columns of metrics
    metrics = [
        ("BFS shells", "depth 2 / 3 / 4", "tested"),
        ("Burned", "39.4% → 30.3%", "(-23% relative)"),
        ("Best individual case", "1 vertex → 12", "(12× ratio)"),
    ]
    for i, (head, val, sub) in enumerate(metrics):
        x = 0.95 + i * 4.0
        add_text(s, x, 5.8, 3.8, 0.25, head,
                 size=10, color="text_subtle", bold=True)
        add_text(s, x, 6.05, 3.8, 0.35, val,
                 size=15, color="green" if "23%" in val or "12×" in sub else "text",
                 bold=True, font=FONT_MONO)
        add_text(s, x, 6.4, 3.8, 0.25, sub,
                 size=10.5, color="text_dim")

    add_text(s, 0.6, 6.95, 12, 0.3,
             "feedback → reformulation → measurable improvement.",
             size=12, color="amber", italic=True, align="center")

    add_page_number(s, 10)


def build_s11_synthetic(prs):
    """Synthetic benchmark ranking (4752 sims)."""
    s = add_slide(prs)
    add_section_badge(s, "RESULTS", color="green")
    add_title(s, "Synthetic benchmark — 4,752 simulations")

    # ranking table — measured on 4,752 synthetic Delaunay runs at k=2
    rows = [
        ("1", "one_step_lookahead",     "27.5", "2.2 ms",   "rose"),
        ("2", "betweenness_front",      "38.5", "17.6 ms",  "purple"),
        ("3", "max_white_neighbors",    "41.5", "0.3 ms",   "blue"),
        ("4", "hybrid_density_aware",   "42.2", "36.0 ms",  "teal"),
        ("5", "max_degree",             "42.4", "0.2 ms",   "blue"),
        ("6", "min_damage_cut",         "42.6", "50.3 ms",  "purple"),
        ("7", "min_cut_edge_front",     "43.4", "11.1 ms",  "purple"),
        ("8", "random (baseline)",      "70.1", "0.2 ms",   "text_subtle"),
    ]
    # header
    cols_x = [0.7, 1.4, 5.6, 7.0]
    cols_w = [0.6, 4.0, 1.3, 1.5]
    y = 2.3
    add_rect(s, cols_x[0], y, sum(cols_w) + 0.4, 0.5,
             fill="card_bg2", rounded=True, corner=0.05)
    for j, h in enumerate(["#", "Strategy", "Burned %", "Cost"]):
        add_text(s, cols_x[j] + 0.05, y + 0.05, cols_w[j], 0.4, h,
                 size=11, color="amber", bold=True,
                 align="left" if j == 1 else "center", anchor="middle")

    for i, (rank, name, pct, cost, col) in enumerate(rows):
        yy = y + 0.55 + i * 0.42
        bg = "card_bg" if i % 2 == 0 else "card_bg2"
        if name.startswith("one_step"):
            add_rect(s, cols_x[0] - 0.05, yy - 0.02, sum(cols_w) + 0.5, 0.42,
                     fill="card_bg", line="amber", line_width=1.4, rounded=True, corner=0.05)
        else:
            add_rect(s, cols_x[0], yy, sum(cols_w) + 0.4, 0.4,
                     fill=bg, rounded=True, corner=0.05)
        add_text(s, cols_x[0], yy, cols_w[0], 0.4, rank,
                 size=12, color="text_dim", align="center", anchor="middle", bold=(rank=="1"))
        add_text(s, cols_x[1], yy, cols_w[1], 0.4, name,
                 size=12, color=col, font=FONT_MONO, anchor="middle",
                 bold=(rank=="1"))
        add_text(s, cols_x[2], yy, cols_w[2], 0.4, pct + "%" if pct != "…" else pct,
                 size=12.5, color="text", bold=(rank in ("1","13")), align="center", anchor="middle", font=FONT_MONO)
        add_text(s, cols_x[3], yy, cols_w[3], 0.4, cost,
                 size=11.5, color="text_dim", align="center", anchor="middle", font=FONT_MONO)

    # right side: takeaways
    add_card(s, 9.4, 2.3, 3.4, 4.5, fill="card_bg2", accent_color="green", accent_side="top")
    add_text(s, 9.6, 2.45, 3.2, 0.35, "TAKEAWAYS",
             size=11, color="green", bold=True)
    takeaways = [
        ("#1 by burn%", "one_step_lookahead", "amber"),
        ("Margin over greedy", "33.8% relative", "purple"),
        ("Random gap", "43 pts (27→70%)", "rose"),
        ("#1 at every size", "n=30, 60, 100", "text_dim"),
    ]
    for i, (head, val, col) in enumerate(takeaways):
        yy = 2.95 + i * 0.85
        add_text(s, 9.6, yy, 3.0, 0.25, head,
                 size=10, color="text_subtle", bold=True)
        add_text(s, 9.6, yy + 0.27, 3.0, 0.35, val,
                 size=12.5, color=col, bold=True, font=FONT_MONO)

    add_text(s, 0.6, 6.95, 12, 0.3,
             "Delaunay-planar synthetic graphs · 3 sizes {30,60,100} × 18 seeds × 11 starts × k=2 × 8 strategies = 4,752 runs.",
             size=10.5, color="text_subtle", italic=True, align="center")

    add_page_number(s, 11)


def build_s12_real_map(prs):
    """Köyceğiz real-map validation."""
    s = add_slide(prs)
    add_section_badge(s, "RESULTS", color="green")
    add_title(s, "Real-map validation — Köyceğiz")

    # stats row
    stats = [
        ("Vertices", "101"),
        ("Active edges", "172"),
        ("Blocked edges", "106"),
        ("Forest total", "50,280 ha"),
        ("Fire origin", "v0 (deg 7)"),
        ("k", "2"),
    ]
    sw = (12.2 / 6) - 0.1; sx = 0.6
    for i, (lab, val) in enumerate(stats):
        x = sx + i * (sw + 0.1)
        add_rect(s, x, 2.3, sw, 0.7, fill="card_bg", line="card_border", line_width=0.5, rounded=True, corner=0.1)
        add_text(s, x, 2.32, sw, 0.25, lab,
                 size=9.5, color="text_subtle", bold=True, align="center", anchor="middle")
        add_text(s, x, 2.55, sw, 0.35, val,
                 size=14, color="amber", bold=True, align="center", anchor="middle", font=FONT_MONO)

    # bar chart (burn % per strategy)
    bars = [
        ("full_rollout_pairs",   15.8, "rose"),
        ("deep_lookahead_3",     15.8, "rose"),
        ("one_step_lookahead",   16.8, "rose"),
        ("articulation_priority",18.8, "purple"),
        ("min_damage_cut",       19.8, "purple"),
        ("max_white_neighbors",  20.8, "blue"),
        ("betweenness_front",    21.5, "purple"),
        ("min_cut_edge_front",   25.0, "purple"),
        ("hybrid_density_aware", 25.5, "teal"),
        ("max_degree",           28.0, "blue"),
        ("random (baseline)",    40.0, "text_subtle"),
    ]
    label_x = 0.7; label_w = 3.4
    bar_x = label_x + label_w + 0.15
    bar_w_max = 7.8
    pct_x = bar_x + bar_w_max + 0.1
    chart_y = 3.3
    row_h = 0.32
    max_pct = max(b[1] for b in bars)
    for i, (name, pct, col) in enumerate(bars):
        y = chart_y + i * row_h
        add_text(s, label_x, y, label_w, row_h, name,
                 size=10.5, color="text_dim", font=FONT_MONO, anchor="middle")
        # bar
        bar_w = bar_w_max * pct / (max_pct * 1.05)
        add_rect(s, bar_x, y + 0.06, bar_w, row_h - 0.1, fill=col, rounded=True, corner=0.2)
        # pct text
        add_text(s, pct_x, y, 0.9, row_h, f"{pct:.1f}%",
                 size=11, color="text", bold=(pct == min(b[1] for b in bars)),
                 anchor="middle", font=FONT_MONO)

    # bottom takeaway
    add_card(s, 0.6, 7.05, 12.2, 0.45, fill="card_bg2", accent_color="green", accent_side="left")
    add_text(s, 0.85, 7.08, 11.7, 0.4,
             "Synthetic ranking confirmed on the real map: lookahead > greedy > random. Deployment validation.",
             size=12, color="text", italic=True, anchor="middle")

    add_page_number(s, 12)


def build_s13_demo(prs):
    s = add_slide(prs)
    add_section_badge(s, "LIVE", color="fire")
    add_title(s, "Live demo — Web Suite")

    # Left: 3-step list
    add_text(s, 0.6, 2.4, 7.0, 0.4, "3-CLICK RUN  ON  STAGE",
             size=12, color="fire", bold=True)
    steps = [
        ("1", "Köyceğiz preset → Build graph",
         "~1 s; in-memory cache warm",
         "fire"),
        ("2", "Click v0 (hub) → Run all strategies",
         "13 strategies in parallel; ~2 s",
         "fire"),
        ("3", "Results → read the recommendation card",
         "graph-specific numbers, not a template",
         "amber"),
    ]
    for i, (num, head, sub, col) in enumerate(steps):
        y = 2.95 + i * 1.05
        # number circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(0.7), Inches(y), Inches(0.55), Inches(0.55))
        circ.fill.solid(); circ.fill.fore_color.rgb = color_of(col); circ.line.fill.background()
        add_text(s, 0.7, y, 0.55, 0.55, num,
                 size=20, color="text", bold=True, align="center", anchor="middle")
        add_text(s, 1.42, y - 0.02, 6.0, 0.4, head,
                 size=15, color="text", bold=True)
        add_text(s, 1.42, y + 0.42, 6.0, 0.4, sub,
                 size=11, color="text_dim", italic=True)

    # Right: large "QR placeholder" + URL
    add_card(s, 8.6, 2.4, 4.2, 4.3, fill="card_bg",
             accent_color="fire", accent_side="top")
    # placeholder QR — checkerboard pattern
    qr_x = 9.3; qr_y = 2.7; qr_size = 2.8
    add_rect(s, qr_x, qr_y, qr_size, qr_size, fill="text", rounded=False)
    add_rect(s, qr_x + 0.15, qr_y + 0.15, qr_size - 0.3, qr_size - 0.3, fill="bg", rounded=False)
    # pseudo modules
    import random
    random.seed(11)
    mod = qr_size / 24
    for r in range(24):
        for cc in range(24):
            if 1 <= r <= 22 and 1 <= cc <= 22 and random.random() > 0.55:
                add_rect(s, qr_x + cc * mod, qr_y + r * mod, mod, mod,
                         fill="text", rounded=False)
    # finder squares
    for (fr, fc) in [(1,1), (1,18), (18,1)]:
        add_rect(s, qr_x + fc * mod, qr_y + fr * mod, mod * 5, mod * 5, fill="text", rounded=False)
        add_rect(s, qr_x + (fc+1) * mod, qr_y + (fr+1) * mod, mod * 3, mod * 3, fill="bg", rounded=False)
        add_rect(s, qr_x + (fc+2) * mod, qr_y + (fr+2) * mod, mod, mod, fill="text", rounded=False)

    add_text(s, 8.6, 5.7, 4.2, 0.4, "http://localhost:5173/",
             size=14, color="amber", bold=True, align="center", font=FONT_MONO)
    add_text(s, 8.6, 6.1, 4.2, 0.35,
             "./run.sh  →  Köyceğiz preset",
             size=10.5, color="text_dim", align="center", font=FONT_MONO)

    # bottom line
    add_text(s, 0.6, 6.95, 12, 0.3,
             "Fallback: identical 90-second demo recording (demo-backup.mp4) ready on USB.",
             size=10, color="text_subtle", italic=True, align="center")

    add_page_number(s, 13)


def build_s14_recommendation(prs):
    s = add_slide(prs)
    add_section_badge(s, "DECISION SUPPORT", color="teal")
    add_title(s, "Recommendation engine — decision, not report")

    # pipeline boxes top
    boxes = [
        ("Graph + fire_origin + k", "INPUT", "blue"),
        ("Topology fingerprint", "n · art · bridge · deg · front", "purple"),
        ("6-rule decision tree", "primary + runner-up", "rose"),
        ("Procedural reason", "graph-specific sentences", "amber"),
    ]
    bw = 2.85; gap = 0.18
    start_x = (13.333 - 4 * bw - 3 * gap) / 2
    by = 2.4
    for i, (head, sub, col) in enumerate(boxes):
        x = start_x + i * (bw + gap)
        rgb = color_of(col)
        fill_rgb = RGBColor(
            int(0.10 * rgb[0] + 0.90 * C["bg"][0]),
            int(0.10 * rgb[1] + 0.90 * C["bg"][1]),
            int(0.10 * rgb[2] + 0.90 * C["bg"][2]),
        )
        add_rect(s, x, by, bw, 1.3, fill=fill_rgb, line=col, line_width=0.8, rounded=True, corner=0.05)
        add_text(s, x + 0.15, by + 0.15, bw - 0.3, 0.4, head,
                 size=13, color="text", bold=True, anchor="middle")
        add_text(s, x + 0.15, by + 0.65, bw - 0.3, 0.5, sub,
                 size=10.5, color=col, anchor="middle", line_spacing=1.3)
        if i < 3:
            ax = x + bw + 0.02
            t = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(ax), Inches(by + 0.58), Inches(0.13), Inches(0.14))
            t.fill.solid(); t.fill.fore_color.rgb = C["amber"]; t.line.fill.background()

    # live-output example
    add_card(s, 0.6, 4.2, 12.2, 2.4, fill="card_bg",
             accent_color="amber", accent_side="left")
    add_text(s, 0.95, 4.35, 5.0, 0.35, "LIVE  OUTPUT  —  v0,  k=2",
             size=11, color="amber", bold=True)
    add_text(s, 0.95, 4.7, 11.4, 0.45, "primary → min_damage_cut",
             size=16, color="green", bold=True, font=FONT_MONO)
    reason = (
        "\"v0 is a hub (7 neighbors, average 3.4). At k=2, five neighbors stay unprotected "
        "this turn. The WHITE subgraph has 7 articulation points + 10 bridges. Min-damage-cut "
        "tests shells at depths 2/3/4 and picks the bottleneck with the best saved/cost ratio "
        "— there are 17 candidate bottlenecks here.\""
    )
    add_text(s, 0.95, 5.2, 11.4, 1.3, reason,
             size=12.5, color="text", italic=True, line_spacing=1.5)

    add_text(s, 0.6, 6.85, 12, 0.3,
             "Designed for OGM operational teams — explains the why, not just the what.",
             size=11, color="text_dim", italic=True, align="center")

    add_page_number(s, 14)


def build_s15_conclusions(prs):
    s = add_slide(prs)
    add_section_badge(s, "CLOSING", color="amber")
    add_title(s, "Conclusions & future work")

    cols = [
        ("Validated  ✓", "green", [
            "4 → 13 strategies (3 philosophies)",
            "Real-map deployment (Köyceğiz)",
            "min_damage_cut answers supervisor's critique",
            "Working web suite delivered",
        ]),
        ("Findings", "amber", [
            "No single dominator — topology + k drive choice",
            "Lookahead wins at scale; betweenness 2nd on cost",
            "Reframing (min-cut → min-damage) measurable",
        ]),
        ("Future", "purple", [
            "Stochastic spread (wind / elevation / fuel)",
            "Multi-period budget (carry-over k)",
            "OGM live integration (real-time fire feed)",
        ]),
    ]
    cw = 4.05; cy = 2.35; ch = 4.2
    cx0 = (13.333 - 3 * cw - 0.3) / 2
    for i, (head, col, items) in enumerate(cols):
        x = cx0 + i * (cw + 0.15)
        rgb = color_of(col)
        fill_rgb = RGBColor(
            int(0.10 * rgb[0] + 0.90 * C["bg"][0]),
            int(0.10 * rgb[1] + 0.90 * C["bg"][1]),
            int(0.10 * rgb[2] + 0.90 * C["bg"][2]),
        )
        add_rect(s, x, cy, cw, ch, fill=fill_rgb, line=col, line_width=0.8, rounded=True, corner=0.04)
        add_text(s, x + 0.25, cy + 0.2, cw - 0.5, 0.5, head,
                 size=18, color=col, bold=True)
        add_line(s, x + 0.25, cy + 0.78, x + cw - 0.25, cy + 0.78, color=col, width=0.8)
        for j, item in enumerate(items):
            yy = cy + 1.0 + j * 0.62
            # bullet
            add_rect(s, x + 0.27, yy + 0.13, 0.10, 0.10, fill=col, rounded=True, corner=0.5)
            add_text(s, x + 0.50, yy, cw - 0.7, 0.6, item,
                     size=11.5, color="text", line_spacing=1.35)

    # Thank you banner
    add_card(s, 4.0, 6.75, 5.3, 0.7, fill="card_bg2", accent_color="amber", accent_side="top")
    add_text(s, 4.0, 6.78, 5.3, 0.65, "Thank you  —  Questions?",
             size=20, color="amber", bold=True, align="center", anchor="middle")

    add_page_number(s, 15)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        build_s1_title,
        build_s2_motivation,
        build_s3_problem,
        build_s4_evolution,
        build_s5_pipeline,
        build_s6_map_to_graph,
        build_s7_edge_filter,
        build_s8_portfolio,
        build_s9_lookahead,
        build_s10_critique,
        build_s11_synthetic,
        build_s12_real_map,
        build_s13_demo,
        build_s14_recommendation,
        build_s15_conclusions,
    ]
    for build in builders:
        build(prs)

    out = Path(__file__).resolve().parent.parent / "SUNUM.pptx"
    prs.save(out)
    print(f"wrote {out}  ({out.stat().st_size / 1024:.0f} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
